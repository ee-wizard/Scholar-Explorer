#!/usr/bin/env python3
"""
Template Generator for md-to-pptx

Generates built-in PPTX template files with predefined layouts.
Templates include: business, tech_dark, education

Each template contains slide layouts:
- Title slide
- Content slide
- Two-column slide
- Image/Chart slide
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


# Template configurations
TEMPLATES = {
    "business": {
        "name": "Business Professional",
        "background": (255, 255, 255),
        "title_color": (44, 62, 80),
        "text_color": (52, 73, 94),
        "accent": (41, 128, 185),
        "secondary": (236, 240, 241),
    },
    "tech_dark": {
        "name": "Tech Dark",
        "background": (30, 30, 30),
        "title_color": (255, 255, 255),
        "text_color": (220, 220, 220),
        "accent": (0, 200, 150),
        "secondary": (45, 45, 45),
    },
    "education": {
        "name": "Education Bright",
        "background": (255, 250, 240),
        "title_color": (70, 130, 180),
        "text_color": (60, 60, 60),
        "accent": (255, 140, 0),
        "secondary": (245, 245, 245),
    },
}


def create_template(config: dict, output_path: str):
    """Create a PPTX template with predefined layouts."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create sample slides demonstrating each layout
    _create_title_layout(prs, config)
    _create_content_layout(prs, config)
    _create_two_column_layout(prs, config)
    _create_image_layout(prs, config)

    prs.save(output_path)
    print(f"  Created: {output_path}")


def _set_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def _create_title_layout(prs, config):
    """Create title slide layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, config["background"])

    # Accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*config["accent"])
    bar.line.fill.background()

    # Title placeholder
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Title Slide]"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*config["title_color"])
    p.alignment = PP_ALIGN.CENTER

    # Subtitle placeholder
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        Inches(12.333), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Subtitle]"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*config["text_color"])
    p.alignment = PP_ALIGN.CENTER


def _create_content_layout(prs, config):
    """Create content slide layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, config["background"])

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*config["accent"])
    bar.line.fill.background()

    # Title area
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Content Title]"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*config["title_color"])

    # Content area
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3),
        Inches(12.333), Inches(5.7)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Content Area - Bullet points, paragraphs, etc.]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(*config["text_color"])


def _create_two_column_layout(prs, config):
    """Create two-column slide layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, config["background"])

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*config["accent"])
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Two Column Title]"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*config["title_color"])

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3),
        Inches(5.9), Inches(5.7)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Left Column]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(*config["text_color"])

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.55), Inches(1.3),
        Inches(0.02), Inches(5.7)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(*config["secondary"])
    divider.line.fill.background()

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(6.9), Inches(1.3),
        Inches(5.9), Inches(5.7)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Right Column]"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(*config["text_color"])


def _create_image_layout(prs, config):
    """Create image/chart focused slide layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, config["background"])

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*config["accent"])
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Image/Chart Title]"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*config["title_color"])

    # Image placeholder area
    img_area = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(1.5),
        Inches(10.333), Inches(5)
    )
    img_area.fill.solid()
    img_area.fill.fore_color.rgb = RGBColor(*config["secondary"])
    img_area.line.color.rgb = RGBColor(*config["accent"])

    # Placeholder text
    placeholder = slide.shapes.add_textbox(
        Inches(4), Inches(3.5),
        Inches(5.333), Inches(1)
    )
    tf = placeholder.text_frame
    p = tf.paragraphs[0]
    p.text = "[Image or Chart Area]"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*config["text_color"])
    p.alignment = PP_ALIGN.CENTER


def generate_all_templates(output_dir: str):
    """Generate all built-in templates."""
    os.makedirs(output_dir, exist_ok=True)

    print(f"Generating templates in: {output_dir}")
    print("-" * 40)

    for name, config in TEMPLATES.items():
        output_path = os.path.join(output_dir, f"{name}_template.pptx")
        create_template(config, output_path)

    print("-" * 40)
    print(f"Generated {len(TEMPLATES)} templates")


def main():
    """Main entry point."""
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate built-in PPTX templates for md-to-pptx"
    )
    parser.add_argument(
        "-o", "--output",
        default=".",
        help="Output directory for templates (default: current directory)"
    )
    parser.add_argument(
        "-t", "--template",
        choices=list(TEMPLATES.keys()),
        help="Generate only specific template"
    )

    args = parser.parse_args()

    if args.template:
        config = TEMPLATES[args.template]
        output_path = os.path.join(args.output, f"{args.template}_template.pptx")
        os.makedirs(args.output, exist_ok=True)
        create_template(config, output_path)
    else:
        generate_all_templates(args.output)


if __name__ == "__main__":
    main()
