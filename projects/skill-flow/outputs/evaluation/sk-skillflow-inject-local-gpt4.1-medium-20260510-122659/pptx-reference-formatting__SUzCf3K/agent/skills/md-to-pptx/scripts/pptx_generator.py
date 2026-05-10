#!/usr/bin/env python3
"""
PPTX Generator Module for md-to-pptx

Generates PowerPoint presentations from structured slide data.
Supports: text, lists, code blocks, tables, images, and charts.
"""

import os
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from enum import Enum

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
except ImportError:
    raise ImportError("python-pptx is required. Install with: pip install python-pptx")

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# Import from md_parser (same directory)
from md_parser import SlideData, ContentElement, ContentType


class ThemeStyle(Enum):
    """Available theme styles."""
    BUSINESS = "business"
    TECH_DARK = "tech_dark"
    EDUCATION = "education"
    NEUMORPHISM = "neumorphism"


@dataclass
class ThemeConfig:
    """Theme configuration for styling."""
    name: str
    background_color: Tuple[int, int, int]
    title_color: Tuple[int, int, int]
    text_color: Tuple[int, int, int]
    accent_color: Tuple[int, int, int]
    code_bg_color: Tuple[int, int, int]
    title_font: str
    body_font: str
    title_size: int  # in points
    body_size: int
    code_size: int


# Built-in themes
THEMES: Dict[str, ThemeConfig] = {
    "business": ThemeConfig(
        name="Business",
        background_color=(255, 255, 255),
        title_color=(44, 62, 80),
        text_color=(52, 73, 94),
        accent_color=(41, 128, 185),
        code_bg_color=(236, 240, 241),
        title_font="Arial",
        body_font="Arial",
        title_size=36,
        body_size=18,
        code_size=14
    ),
    "tech_dark": ThemeConfig(
        name="Tech Dark",
        background_color=(30, 30, 30),
        title_color=(255, 255, 255),
        text_color=(220, 220, 220),
        accent_color=(0, 200, 150),
        code_bg_color=(45, 45, 45),
        title_font="Consolas",
        body_font="Segoe UI",
        title_size=36,
        body_size=18,
        code_size=14
    ),
    "education": ThemeConfig(
        name="Education",
        background_color=(255, 250, 240),
        title_color=(70, 130, 180),
        text_color=(60, 60, 60),
        accent_color=(255, 140, 0),
        code_bg_color=(245, 245, 245),
        title_font="Georgia",
        body_font="Verdana",
        title_size=36,
        body_size=18,
        code_size=14
    ),
    "neumorphism": ThemeConfig(
        name="Neumorphism",
        background_color=(240, 243, 249),
        title_color=(45, 55, 72),
        text_color=(74, 85, 104),
        accent_color=(66, 153, 225),
        code_bg_color=(226, 232, 240),
        title_font="Arial",
        body_font="Arial",
        title_size=36,
        body_size=18,
        code_size=14
    )
}


class PPTXGenerator:
    """Generate PPTX files from slide data."""

    # Slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Content margins
    MARGIN_LEFT = Inches(0.5)
    MARGIN_RIGHT = Inches(0.5)
    MARGIN_TOP = Inches(1.2)
    MARGIN_BOTTOM = Inches(0.5)

    def __init__(self, theme: str = "business", template_path: Optional[str] = None):
        """
        Initialize generator with theme or template.

        Args:
            theme: Theme name (business, tech_dark, education, neumorphism)
            template_path: Optional path to custom .pptx template
        """
        self.template_path = template_path
        self.theme = THEMES.get(theme, THEMES["business"])
        self.prs: Optional[Presentation] = None
        self.warnings: List[str] = []
        self.progress_callback = None
        self._use_template_background = template_path is not None

    def set_progress_callback(self, callback):
        """Set callback function for progress updates."""
        self.progress_callback = callback

    def _report_progress(self, current: int, total: int, message: str):
        """Report progress if callback is set."""
        if self.progress_callback:
            self.progress_callback(current, total, message)

    def generate(self, slides: List[SlideData], output_path: str) -> bool:
        """
        Generate PPTX from slide data.

        Args:
            slides: List of SlideData objects
            output_path: Output file path

        Returns:
            True if successful
        """
        self.warnings = []

        # Create presentation
        if self.template_path and os.path.exists(self.template_path):
            self.prs = Presentation(self.template_path)
        else:
            self.prs = Presentation()
            self.prs.slide_width = self.SLIDE_WIDTH
            self.prs.slide_height = self.SLIDE_HEIGHT

        total_slides = len(slides)

        for i, slide_data in enumerate(slides):
            self._report_progress(i + 1, total_slides, f"Generating slide {i + 1}: {slide_data.title[:30]}...")

            try:
                if slide_data.layout_hint == "title":
                    self._create_title_slide(slide_data)
                elif slide_data.layout_hint == "image":
                    self._create_image_slide(slide_data)
                elif slide_data.layout_hint == "chart":
                    self._create_chart_slide(slide_data)
                elif slide_data.layout_hint == "two_column":
                    self._create_two_column_slide(slide_data)
                else:
                    self._create_content_slide(slide_data)
            except Exception as e:
                self.warnings.append(f"Error on slide {i + 1}: {str(e)}")
                # Create fallback slide
                self._create_fallback_slide(slide_data, str(e))

        # Save presentation
        try:
            self.prs.save(output_path)
            self._report_progress(total_slides, total_slides, f"Saved to {output_path}")
            return True
        except Exception as e:
            self.warnings.append(f"Failed to save: {str(e)}")
            return False

    def _add_blank_slide(self) -> Any:
        """Add a blank slide to the presentation."""
        # Try to find a blank layout, fallback to last available layout
        layouts = self.prs.slide_layouts
        num_layouts = len(layouts)

        # Common blank layout indices to try
        blank_indices = [6, 5, num_layouts - 1, 0]

        for idx in blank_indices:
            if idx < num_layouts:
                try:
                    return self.prs.slides.add_slide(layouts[idx])
                except Exception:
                    continue

        # Last resort: use first available layout
        return self.prs.slides.add_slide(layouts[0])

    def _set_background(self, slide):
        """Set slide background color based on theme (skip if using template)."""
        if self._use_template_background:
            # Don't override template background
            return
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self.theme.background_color)
        except Exception:
            # Some templates may not support background modification
            pass

    def _add_title_shape(self, slide, title: str, top: float = 0.3, height: float = 0.8) -> Any:
        """Add title text box to slide."""
        left = self.MARGIN_LEFT
        width = self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT

        shape = slide.shapes.add_textbox(left, Inches(top), width, Inches(height))
        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.theme.title_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.title_color)
        p.font.name = self.theme.title_font
        p.alignment = PP_ALIGN.LEFT

        return shape

    def _create_title_slide(self, slide_data: SlideData):
        """Create a title slide."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Center title
        title_top = 2.5
        title_shape = slide.shapes.add_textbox(
            self.MARGIN_LEFT,
            Inches(title_top),
            self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT,
            Inches(1.5)
        )
        tf = title_shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.title_color)
        p.font.name = self.theme.title_font
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if slide_data.subtitle:
            subtitle_shape = slide.shapes.add_textbox(
                self.MARGIN_LEFT,
                Inches(title_top + 1.5),
                self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT,
                Inches(1)
            )
            tf = subtitle_shape.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(*self.theme.text_color)
            p.font.name = self.theme.body_font
            p.alignment = PP_ALIGN.CENTER

    def _create_content_slide(self, slide_data: SlideData):
        """Create a standard content slide."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Add title
        self._add_title_shape(slide, slide_data.title)

        # Content area
        content_top = self.MARGIN_TOP
        content_left = self.MARGIN_LEFT
        content_width = self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT
        available_height = self.SLIDE_HEIGHT - self.MARGIN_TOP - self.MARGIN_BOTTOM

        current_top = content_top

        for element in slide_data.elements:
            if current_top > self.SLIDE_HEIGHT - self.MARGIN_BOTTOM:
                self.warnings.append(f"Content overflow on slide: {slide_data.title}")
                break

            element_height = self._add_element(
                slide, element, content_left, current_top, content_width
            )
            current_top += element_height + Inches(0.2)

    def _create_image_slide(self, slide_data: SlideData):
        """Create a slide with image focus."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Add title
        self._add_title_shape(slide, slide_data.title)

        # Find image element
        image_elem = None
        other_elements = []
        for elem in slide_data.elements:
            if elem.type == ContentType.IMAGE and image_elem is None:
                image_elem = elem
            else:
                other_elements.append(elem)

        if image_elem:
            # Add image centered
            self._add_image_element(
                slide, image_elem,
                Inches(2), Inches(1.5),
                self.SLIDE_WIDTH - Inches(4)
            )

        # Add other elements below
        current_top = Inches(5.5)
        for elem in other_elements[:2]:  # Limit to 2 extra elements
            height = self._add_element(
                slide, elem,
                self.MARGIN_LEFT, current_top,
                self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT
            )
            current_top += height + Inches(0.1)

    def _create_chart_slide(self, slide_data: SlideData):
        """Create a slide with chart focus."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Add title
        self._add_title_shape(slide, slide_data.title)

        # Find chart or table data
        chart_elem = None
        for elem in slide_data.elements:
            if elem.type in (ContentType.CHART_DATA, ContentType.TABLE):
                chart_elem = elem
                break

        if chart_elem:
            if chart_elem.type == ContentType.TABLE:
                # Convert table to chart
                self._add_chart_from_table(
                    slide, chart_elem,
                    Inches(1), Inches(1.5),
                    self.SLIDE_WIDTH - Inches(2),
                    Inches(5)
                )
            else:
                self._add_chart_element(
                    slide, chart_elem,
                    Inches(1), Inches(1.5),
                    self.SLIDE_WIDTH - Inches(2),
                    Inches(5)
                )

    def _create_two_column_slide(self, slide_data: SlideData):
        """Create a two-column layout slide."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Add title
        self._add_title_shape(slide, slide_data.title)

        # Split elements into two columns
        elements = slide_data.elements
        mid = len(elements) // 2

        left_elements = elements[:mid]
        right_elements = elements[mid:]

        col_width = (self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT - Inches(0.5)) / 2

        # Left column
        current_top = self.MARGIN_TOP
        for elem in left_elements:
            height = self._add_element(slide, elem, self.MARGIN_LEFT, current_top, col_width)
            current_top += height + Inches(0.15)

        # Right column
        current_top = self.MARGIN_TOP
        right_left = self.MARGIN_LEFT + col_width + Inches(0.5)
        for elem in right_elements:
            height = self._add_element(slide, elem, right_left, current_top, col_width)
            current_top += height + Inches(0.15)

    def _create_fallback_slide(self, slide_data: SlideData, error: str):
        """Create a simple fallback slide when normal creation fails."""
        slide = self._add_blank_slide()
        self._set_background(slide)

        # Add title
        self._add_title_shape(slide, slide_data.title)

        # Add error note
        shape = slide.shapes.add_textbox(
            self.MARGIN_LEFT,
            self.MARGIN_TOP,
            self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT,
            Inches(1)
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Content generation error: {error}]"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 100, 100)

    def _add_element(self, slide, element: ContentElement, left, top, width) -> float:
        """
        Add a content element to the slide.

        Returns:
            Height of the added element in EMUs
        """
        if element.type == ContentType.PARAGRAPH:
            return self._add_paragraph_element(slide, element, left, top, width)
        elif element.type == ContentType.HEADING:
            return self._add_heading_element(slide, element, left, top, width)
        elif element.type in (ContentType.BULLET_LIST, ContentType.NUMBERED_LIST):
            return self._add_list_element(slide, element, left, top, width)
        elif element.type == ContentType.CODE_BLOCK:
            return self._add_code_element(slide, element, left, top, width)
        elif element.type == ContentType.TABLE:
            return self._add_table_element(slide, element, left, top, width)
        elif element.type == ContentType.IMAGE:
            return self._add_image_element(slide, element, left, top, width)
        else:
            return Inches(0)

    def _add_paragraph_element(self, slide, element: ContentElement, left, top, width) -> float:
        """Add paragraph text."""
        height = Inches(0.5)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = element.content
        p.font.size = Pt(self.theme.body_size)
        p.font.color.rgb = RGBColor(*self.theme.text_color)
        p.font.name = self.theme.body_font

        return height

    def _add_heading_element(self, slide, element: ContentElement, left, top, width) -> float:
        """Add sub-heading text."""
        level = element.metadata.get("level", 3)
        size = max(14, self.theme.body_size + (6 - level) * 2)

        height = Inches(0.4)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame

        p = tf.paragraphs[0]
        p.text = element.content
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.theme.title_color)
        p.font.name = self.theme.title_font

        return height

    def _add_list_element(self, slide, element: ContentElement, left, top, width) -> float:
        """Add bullet or numbered list."""
        items = element.content
        is_numbered = element.type == ContentType.NUMBERED_LIST

        line_height = Inches(0.35)
        height = line_height * len(items)

        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            indent = item.get("indent", 0)
            prefix = f"{i + 1}. " if is_numbered else "• "
            indent_str = "    " * indent

            p.text = f"{indent_str}{prefix}{item['text']}"
            p.font.size = Pt(self.theme.body_size)
            p.font.color.rgb = RGBColor(*self.theme.text_color)
            p.font.name = self.theme.body_font
            p.space_after = Pt(6)

        return height

    def _add_code_element(self, slide, element: ContentElement, left, top, width) -> float:
        """Add code block with simple styling."""
        code = element.content
        lines = code.split('\n')
        line_count = min(len(lines), 15)  # Limit lines

        line_height = Inches(0.25)
        height = line_height * line_count + Inches(0.3)

        # Background shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(*self.theme.code_bg_color)
        bg_shape.line.fill.background()

        # Code text
        text_shape = slide.shapes.add_textbox(
            left + Inches(0.1),
            top + Inches(0.1),
            width - Inches(0.2),
            height - Inches(0.2)
        )
        tf = text_shape.text_frame
        tf.word_wrap = False

        display_code = '\n'.join(lines[:line_count])
        if len(lines) > line_count:
            display_code += f"\n... ({len(lines) - line_count} more lines)"

        p = tf.paragraphs[0]
        p.text = display_code
        p.font.size = Pt(self.theme.code_size)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(*self.theme.text_color)

        return height

    def _add_table_element(self, slide, element: ContentElement, left, top, width) -> float:
        """Add table."""
        data = element.content
        headers = data.get("headers", [])
        rows = data.get("rows", [])

        if not headers and not rows:
            return Inches(0)

        col_count = len(headers) if headers else len(rows[0]) if rows else 0
        row_count = len(rows) + (1 if headers else 0)

        if col_count == 0 or row_count == 0:
            return Inches(0)

        row_height = Inches(0.4)
        height = row_height * row_count

        table = slide.shapes.add_table(
            row_count, col_count,
            left, top, width, height
        ).table

        # Style table
        for i, cell in enumerate(table.rows[0].cells if headers else []):
            cell.text = headers[i] if i < len(headers) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self.theme.accent_color)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(self.theme.body_size - 2)
            p.font.color.rgb = RGBColor(255, 255, 255)

        start_row = 1 if headers else 0
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < col_count:
                    cell = table.rows[start_row + row_idx].cells[col_idx]
                    cell.text = str(cell_text)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(self.theme.body_size - 2)
                    p.font.color.rgb = RGBColor(*self.theme.text_color)

        return height

    def _add_image_element(self, slide, element: ContentElement, left, top, width) -> float:
        """Add image from local file."""
        img_path = element.content

        if not os.path.exists(img_path):
            self.warnings.append(f"Image not found: {img_path}")
            return Inches(0)

        try:
            # Get image dimensions
            max_height = Inches(4)

            if HAS_PIL:
                with Image.open(img_path) as img:
                    img_width, img_height = img.size
                    aspect = img_width / img_height

                    # Calculate dimensions to fit
                    if width / aspect > max_height:
                        height = max_height
                        actual_width = height * aspect
                    else:
                        actual_width = width
                        height = width / aspect
            else:
                actual_width = width
                height = max_height

            # Center image
            center_left = left + (width - actual_width) / 2

            slide.shapes.add_picture(img_path, center_left, top, actual_width, height)
            return height

        except Exception as e:
            self.warnings.append(f"Failed to add image {img_path}: {str(e)}")
            return Inches(0)

    def _add_chart_from_table(self, slide, element: ContentElement, left, top, width, height):
        """Create a chart from table data."""
        data = element.content
        headers = data.get("headers", [])
        rows = data.get("rows", [])

        if len(headers) < 2 or not rows:
            # Fall back to table
            self._add_table_element(slide, element, left, top, width)
            return

        try:
            chart_data = CategoryChartData()
            chart_data.categories = [row[0] for row in rows]

            # Add series for each numeric column
            for col_idx in range(1, len(headers)):
                values = []
                for row in rows:
                    if col_idx < len(row):
                        try:
                            values.append(float(row[col_idx].replace(',', '')))
                        except (ValueError, AttributeError):
                            values.append(0)
                    else:
                        values.append(0)
                chart_data.add_series(headers[col_idx], values)

            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                left, top, width, height,
                chart_data
            ).chart

            # Style chart
            chart.has_legend = len(headers) > 2

        except Exception as e:
            self.warnings.append(f"Chart creation failed: {str(e)}")
            self._add_table_element(slide, element, left, top, width)

    def _add_chart_element(self, slide, element: ContentElement, left, top, width, height):
        """Add chart from chart data."""
        # Placeholder for direct chart data support
        pass

    def get_warnings(self) -> List[str]:
        """Return warnings generated during generation."""
        return self.warnings


def generate_pptx(
    slides: List[SlideData],
    output_path: str,
    theme: str = "business",
    template_path: Optional[str] = None,
    progress_callback=None
) -> Tuple[bool, List[str]]:
    """
    Convenience function to generate PPTX.

    Args:
        slides: List of SlideData objects
        output_path: Output file path
        theme: Theme name
        template_path: Optional custom template path
        progress_callback: Optional callback(current, total, message)

    Returns:
        Tuple of (success, warnings)
    """
    generator = PPTXGenerator(theme=theme, template_path=template_path)
    if progress_callback:
        generator.set_progress_callback(progress_callback)

    success = generator.generate(slides, output_path)
    return success, generator.get_warnings()


if __name__ == "__main__":
    # Test with sample data
    from md_parser import parse_markdown

    sample_md = """
# Test Presentation

A demo presentation

## Introduction

Welcome to this test presentation.

- First point
- Second point
- Third point

## Code Example

```python
def hello():
    print("Hello!")
```

## Data

| Item | Value |
|------|-------|
| A    | 100   |
| B    | 200   |
| C    | 150   |

## Thank You

Questions?
"""

    slides, parse_warnings = parse_markdown(sample_md)
    print(f"Parsed {len(slides)} slides")

    def progress(current, total, msg):
        print(f"[{current}/{total}] {msg}")

    success, gen_warnings = generate_pptx(
        slides,
        "test_output.pptx",
        theme="business",
        progress_callback=progress
    )

    print(f"\nGeneration {'succeeded' if success else 'failed'}")
    if parse_warnings:
        print(f"Parse warnings: {parse_warnings}")
    if gen_warnings:
        print(f"Generation warnings: {gen_warnings}")
