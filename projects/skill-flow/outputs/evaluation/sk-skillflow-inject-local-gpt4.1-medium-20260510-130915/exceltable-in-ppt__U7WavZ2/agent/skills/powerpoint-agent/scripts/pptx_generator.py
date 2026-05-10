"""Create PowerPoint presentations from slide images using python-pptx.

Each slide is a full-slide AI-generated image with slide titles stored
in speaker notes for accessibility.
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches


def create_pptx_from_images(
    image_paths: list[str],
    output_path: str,
    slide_titles: Optional[list[str]] = None,
    slide_notes: Optional[list[str]] = None,
) -> str:
    """Create PowerPoint presentation with full-slide images.

    Each slide contains a single AI-generated image stretched to fill
    the entire slide. Slide titles are added to speaker notes for
    accessibility and reference.

    Args:
        image_paths: List of paths to slide images (PNG, JPG, etc.).
        output_path: Path for output .pptx file.
        slide_titles: Optional list of slide titles (added to notes).
        slide_notes: Optional list of additional notes for each slide.

    Returns:
        Path to the created PowerPoint file.

    Raises:
        FileNotFoundError: If any image file doesn't exist.
        ValueError: If image_paths is empty.
    """
    if not image_paths:
        raise ValueError("image_paths cannot be empty")

    # Validate all images exist
    for path in image_paths:
        if not Path(path).exists():
            raise FileNotFoundError(f"Image not found: {path}")

    # Create presentation with 16:9 widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(10)     # Standard widescreen width
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # Use blank layout (index 6) for clean image-only slides
    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(image_paths):
        # Add slide with blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add image stretched to fill entire slide
        slide.shapes.add_picture(
            img_path,
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Add title and notes to speaker notes for accessibility
        notes_text = []

        if slide_titles and i < len(slide_titles):
            notes_text.append(f"Title: {slide_titles[i]}")

        if slide_notes and i < len(slide_notes):
            notes_text.append(f"\n{slide_notes[i]}")

        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "".join(notes_text)

    # Ensure output directory exists
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    # Save presentation
    prs.save(output_path)

    return output_path


def create_pptx_with_metadata(
    image_paths: list[str],
    output_path: str,
    slides_metadata: list[dict],
    presentation_title: Optional[str] = None,
    presentation_author: Optional[str] = None,
) -> str:
    """Create PowerPoint with slide metadata and presentation properties.

    Extended version of create_pptx_from_images that includes
    presentation-level metadata and per-slide metadata.

    Args:
        image_paths: List of paths to slide images.
        output_path: Path for output .pptx file.
        slides_metadata: List of dicts with slide info. Expected keys:
            - title: Slide title
            - bullet_points: List of bullet points (for notes)
            - narration: Speaker notes/script
        presentation_title: Optional presentation title for properties.
        presentation_author: Optional author name for properties.

    Returns:
        Path to the created PowerPoint file.
    """
    if not image_paths:
        raise ValueError("image_paths cannot be empty")

    # Create presentation with 16:9 widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Set presentation properties if provided
    if presentation_title:
        prs.core_properties.title = presentation_title
    if presentation_author:
        prs.core_properties.author = presentation_author

    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(image_paths):
        slide = prs.slides.add_slide(blank_layout)

        # Add full-slide image
        slide.shapes.add_picture(
            img_path,
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Build notes from metadata
        if i < len(slides_metadata):
            meta = slides_metadata[i]
            notes_parts = []

            # Title
            if meta.get("title"):
                notes_parts.append(f"Title: {meta['title']}")

            # Bullet points
            if meta.get("bullet_points"):
                bullets = "\n".join(f"  - {bp}" for bp in meta["bullet_points"])
                notes_parts.append(f"\nKey Points:\n{bullets}")

            # Narration/script
            if meta.get("narration"):
                notes_parts.append(f"\nSpeaker Notes:\n{meta['narration']}")

            if notes_parts:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = "\n".join(notes_parts)

    # Ensure output directory exists
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    prs.save(output_path)

    return output_path


def get_pptx_info(pptx_path: str) -> dict:
    """Get information about a PowerPoint file.

    Args:
        pptx_path: Path to .pptx file.

    Returns:
        Dictionary with presentation info:
            - slide_count: Number of slides
            - width_inches: Slide width
            - height_inches: Slide height
            - title: Presentation title (if set)
            - author: Author (if set)
    """
    prs = Presentation(pptx_path)

    return {
        "slide_count": len(prs.slides),
        "width_inches": prs.slide_width.inches,
        "height_inches": prs.slide_height.inches,
        "aspect_ratio": f"{prs.slide_width.inches / prs.slide_height.inches:.2f}:1",
        "title": prs.core_properties.title or "",
        "author": prs.core_properties.author or "",
    }


def replace_slide_image(
    pptx_path: str,
    slide_number: int,
    new_image_path: str,
    output_path: Optional[str] = None,
    new_title: Optional[str] = None,
    new_notes: Optional[str] = None,
) -> str:
    """Replace the image on a specific slide in an existing PowerPoint.

    Args:
        pptx_path: Path to existing .pptx file.
        slide_number: 1-based slide number to replace (1 = first slide).
        new_image_path: Path to new image file.
        output_path: Path for output file. If None, overwrites original.
        new_title: Optional new title for speaker notes.
        new_notes: Optional new speaker notes content.

    Returns:
        Path to the modified PowerPoint file.

    Raises:
        FileNotFoundError: If pptx or image file doesn't exist.
        ValueError: If slide_number is out of range.
    """
    if not Path(pptx_path).exists():
        raise FileNotFoundError(f"PowerPoint not found: {pptx_path}")
    if not Path(new_image_path).exists():
        raise FileNotFoundError(f"Image not found: {new_image_path}")

    prs = Presentation(pptx_path)

    # Validate slide number (1-based)
    if slide_number < 1 or slide_number > len(prs.slides):
        raise ValueError(
            f"Slide number {slide_number} out of range. "
            f"Presentation has {len(prs.slides)} slides."
        )

    # Get the slide (convert to 0-based index)
    slide = prs.slides[slide_number - 1]

    # Remove existing picture shapes from the slide
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Add new full-slide image
    slide.shapes.add_picture(
        new_image_path,
        left=Inches(0),
        top=Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )

    # Update notes if provided
    if new_title or new_notes:
        notes_parts = []
        if new_title:
            notes_parts.append(f"Title: {new_title}")
        if new_notes:
            notes_parts.append(f"\n{new_notes}")

        if notes_parts:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "".join(notes_parts)

    # Save to output path or overwrite original
    save_path = output_path or pptx_path
    output_dir = Path(save_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    prs.save(save_path)
    return save_path


def replace_multiple_slides(
    pptx_path: str,
    replacements: list[dict],
    output_path: Optional[str] = None,
) -> str:
    """Replace images on multiple slides in an existing PowerPoint.

    Args:
        pptx_path: Path to existing .pptx file.
        replacements: List of dicts with replacement info:
            - slide_number: 1-based slide number
            - image_path: Path to new image
            - title: Optional new title for notes
            - notes: Optional new notes content
        output_path: Path for output file. If None, overwrites original.

    Returns:
        Path to the modified PowerPoint file.
    """
    if not Path(pptx_path).exists():
        raise FileNotFoundError(f"PowerPoint not found: {pptx_path}")

    prs = Presentation(pptx_path)

    for repl in replacements:
        slide_num = repl["slide_number"]
        img_path = repl["image_path"]

        if not Path(img_path).exists():
            raise FileNotFoundError(f"Image not found: {img_path}")

        if slide_num < 1 or slide_num > len(prs.slides):
            raise ValueError(
                f"Slide number {slide_num} out of range. "
                f"Presentation has {len(prs.slides)} slides."
            )

        slide = prs.slides[slide_num - 1]

        # Remove existing pictures
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        # Add new image
        slide.shapes.add_picture(
            img_path,
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Update notes if provided
        if repl.get("title") or repl.get("notes"):
            notes_parts = []
            if repl.get("title"):
                notes_parts.append(f"Title: {repl['title']}")
            if repl.get("notes"):
                notes_parts.append(f"\n{repl['notes']}")

            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "".join(notes_parts)

    # Save
    save_path = output_path or pptx_path
    output_dir = Path(save_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    prs.save(save_path)
    return save_path
